"""Marp(plass 테마) template.md 구조를 python-pptx로 포팅한 생성기.

CSS(flex/grid) 레이아웃을 pptx의 절대 좌표 방식으로 근사 변환한 것이라
Marp 출력과 100% 동일하지는 않음. 실행: python3 scripts/build_pptx.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image
import os

ROOT = os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
ASSETS = os.path.join(ROOT, "assets")
OUT = os.path.join(ROOT, "dist", "template.pptx")

FONT = "Pretendard"  # 미설치 환경에선 시스템이 자동 대체 (예: 맑은 고딕)

PRIMARY = RGBColor(0x1F, 0x47, 0x89)
ACCENT = RGBColor(0xFF, 0x6B, 0x35)
TEXT = RGBColor(0x1D, 0x1D, 0x1D)
TEXT_INV = RGBColor(0xFF, 0xFF, 0xFF)
MUTED = RGBColor(0x77, 0x77, 0x77)  # opacity:0.6 텍스트의 근사값
GRADIENT = [
    (0.0, RGBColor(0x10, 0x23, 0x45)),
    (0.6, RGBColor(0x0C, 0x5C, 0xBC)),
    (1.0, RGBColor(0x58, 0x89, 0xD9)),
]

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN_X = Inches(0.6)

BG1 = os.path.join(ASSETS, "bg1.jpg")
BG2 = os.path.join(ASSETS, "bg2.jpg")


def new_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_fullbleed_image(slide, path):
    slide.shapes.add_picture(path, 0, 0, width=SLIDE_W, height=SLIDE_H)


def add_gradient_bar(slide):
    bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_X, 0, SLIDE_W - 2 * MARGIN_X, Inches(0.15)
    )
    bar.line.fill.background()
    fill = bar.fill
    fill.gradient()
    stops = fill.gradient_stops
    # python-pptx는 기본 2 stop만 제공 -> XML 직접 조작으로 3 stop 구성
    gs_lst = fill._xPr.find(qn("a:gradFill")).find(qn("a:gsLst"))
    for el in list(gs_lst):
        gs_lst.remove(el)
    for pos, color in GRADIENT:
        gs = gs_lst.makeelement(qn("a:gs"), {"pos": str(int(pos * 100000))})
        srgb = gs.makeelement(qn("a:srgbClr"), {"val": "%02X%02X%02X" % (color[0], color[1], color[2])})
        gs.append(srgb)
        gs_lst.append(gs)
    fill.gradient_angle = 0
    return bar


def add_underline(slide, x, y, width, color=PRIMARY, weight=Pt(1)):
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, width, Emu(1))
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    line.height = Emu(int(weight))
    return line


def add_text(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=Pt(18),
    color=TEXT,
    bold=False,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
    font=FONT,
):
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = size
    run.font.bold = bold
    run.font.name = font
    run.font.color.rgb = color
    return box


def add_bullets(slide, x, y, w, h, items, size=Pt(18), color=TEXT, font=FONT):
    """items: [(text, level)] level 0부터 시작 (마크다운 들여쓰기 단계)"""
    box = slide.shapes.add_textbox(x, y, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, (text, level) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.level = level
        p.text = "•  " + text if level == 0 else "-  " + text
        for run in p.runs:
            run.font.size = size
            run.font.name = font
            run.font.color.rgb = color
    return box


def content_header(slide, title, subtitle=None):
    add_gradient_bar(slide)
    title_y = Inches(0.32)
    add_text(
        slide,
        MARGIN_X,
        title_y,
        SLIDE_W - 2 * MARGIN_X,
        Inches(0.55),
        title,
        size=Pt(28),
        color=PRIMARY,
        bold=True,
    )
    add_underline(slide, MARGIN_X, title_y + Inches(0.52), SLIDE_W - 2 * MARGIN_X)
    body_top = title_y + Inches(0.6)
    if subtitle:
        add_text(
            slide,
            MARGIN_X,
            body_top,
            SLIDE_W - 2 * MARGIN_X,
            Inches(0.4),
            subtitle,
            size=Pt(16),
            color=TEXT,
        )
        body_top += Inches(0.45)
    return body_top


def add_footnotes(slide, footnotes):
    if not footnotes:
        return
    y = SLIDE_H - Inches(0.28) * len(footnotes) - Inches(0.15)
    for i, text in enumerate(footnotes):
        add_text(
            slide,
            MARGIN_X,
            y + Inches(0.24) * i,
            SLIDE_W - 2 * MARGIN_X,
            Inches(0.24),
            f"* {text}",
            size=Pt(9),
            color=MUTED,
        )


def add_page_number(slide, current, total):
    add_text(
        slide,
        SLIDE_W - Inches(1.2),
        SLIDE_H - Inches(0.35),
        Inches(1.0),
        Inches(0.3),
        f"[{current}/{total}]",
        size=Pt(9),
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


# ---------------------------------------------------------------------------
# 슬라이드 종류별 빌더
# ---------------------------------------------------------------------------

def add_lead_slide(prs, title, subtitle, meta):
    slide = blank_slide(prs)
    add_fullbleed_image(slide, BG1)

    add_text(
        slide,
        Inches(0.9),
        Inches(2.7),
        Inches(8.0),
        Inches(0.8),
        title,
        size=Pt(34),
        color=PRIMARY,
        bold=True,
    )
    add_underline(slide, Inches(0.95), Inches(3.35), Inches(1.7), weight=Pt(2))
    add_text(
        slide,
        Inches(0.9),
        Inches(3.5),
        Inches(8.0),
        Inches(0.6),
        subtitle,
        size=Pt(18),
        color=TEXT,
    )

    meta_box = slide.shapes.add_textbox(Inches(8.5), Inches(5.6), Inches(4.2), Inches(1.4))
    tf = meta_box.text_frame
    tf.word_wrap = True
    for i, line in enumerate(meta):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = line
        run.font.size = Pt(14)
        run.font.name = FONT
        run.font.color.rgb = TEXT
    return slide


def add_toc_slide(prs, items):
    slide = blank_slide(prs)
    add_fullbleed_image(slide, BG1)

    title_y = Inches(1.0)
    add_underline(slide, Inches(0.9), title_y - Inches(0.15), Inches(1.7), weight=Pt(4))
    add_text(
        slide,
        Inches(0.9),
        title_y,
        Inches(4.0),
        Inches(0.7),
        "Contents",
        size=Pt(34),
        color=PRIMARY,
        bold=True,
    )
    add_underline(slide, Inches(0.9), title_y + Inches(0.85), Inches(1.7), weight=Pt(1))

    y = title_y + Inches(1.3)
    for i, item in enumerate(items, start=1):
        add_text(
            slide,
            Inches(0.9),
            y,
            Inches(0.8),
            Inches(0.5),
            f"{i:02d}",
            size=Pt(18),
            color=PRIMARY,
            bold=True,
        )
        add_text(
            slide,
            Inches(1.6),
            y,
            Inches(9.0),
            Inches(0.5),
            item,
            size=Pt(18),
            color=TEXT,
            bold=True,
        )
        y += Inches(0.6)
    return slide


def add_content_slide(prs, title, subtitle, columns, footnotes=None):
    """columns: N개의 리스트, 각 원소는 [(text, level), ...] (N = 1~4단)"""
    slide = blank_slide(prs)
    body_top = content_header(slide, title, subtitle)

    n = len(columns)
    gap = Inches(0.3)
    total_w = SLIDE_W - 2 * MARGIN_X - gap * (n - 1)
    col_w = Emu(int(total_w / n))
    body_h = SLIDE_H - body_top - Inches(0.6)

    for i, items in enumerate(columns):
        x = MARGIN_X + (col_w + gap) * i
        add_bullets(slide, x, body_top, col_w, body_h, items)

    add_footnotes(slide, footnotes)
    return slide


def add_image_slide(prs, title, subtitle, images, footnotes=None):
    """images: 1~4개의 이미지 경로 리스트 (컬럼 수 = 이미지 수)"""
    slide = blank_slide(prs)
    body_top = content_header(slide, title, subtitle)

    n = len(images)
    gap = Inches(0.3)
    total_w = SLIDE_W - 2 * MARGIN_X - gap * (n - 1)
    col_w = int(total_w / n)
    body_h = int(SLIDE_H - body_top - Inches(0.6))

    for i, path in enumerate(images):
        cell_x = int(MARGIN_X) + (col_w + int(gap)) * i
        cell_y = int(body_top)

        im = Image.open(path)
        iw, ih = im.size
        scale = min(col_w / iw, body_h / ih)
        w, h = int(iw * scale), int(ih * scale)
        x = cell_x + (col_w - w) // 2
        y = cell_y + (body_h - h) // 2
        slide.shapes.add_picture(path, x, y, width=w, height=h)

    add_footnotes(slide, footnotes)
    return slide


def add_table_slide(prs, title, subtitle, headers, rows, footnotes=None):
    slide = blank_slide(prs)
    body_top = content_header(slide, title, subtitle)

    n_rows = len(rows) + 1
    n_cols = len(headers)
    table_w = SLIDE_W - 2 * MARGIN_X
    table_h = min(SLIDE_H - body_top - Inches(0.6), Inches(0.5) * n_rows)

    gframe = slide.shapes.add_table(n_rows, n_cols, MARGIN_X, body_top, table_w, table_h)
    table = gframe.table

    for c, text in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = text
        cell.fill.solid()
        cell.fill.fore_color.rgb = PRIMARY
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                run.font.color.rgb = TEXT_INV
                run.font.bold = True
                run.font.size = Pt(14)
                run.font.name = FONT

    for r, row in enumerate(rows, start=1):
        for c, text in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(text)
            cell.fill.solid()
            cell.fill.fore_color.rgb = RGBColor(0xF0, 0xF3, 0xF9) if r % 2 == 0 else TEXT_INV
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    run.font.color.rgb = TEXT
                    run.font.size = Pt(13)
                    run.font.name = FONT

    add_footnotes(slide, footnotes)
    return slide


# ---------------------------------------------------------------------------
# template.md 데모 구성 그대로 재현
# ---------------------------------------------------------------------------

def build_demo():
    prs = new_presentation()

    add_lead_slide(
        prs,
        "제목을 입력하세요",
        "부제목을 입력하세요",
        ["소속을 입력해세요", "이름을 입력하세요", "발표일을 입력하세요"],
    )

    add_toc_slide(
        prs,
        [f"목차를 입력하세요" for _ in range(5)],
    )

    add_content_slide(
        prs,
        "슬라이드 제목을 입력하세요 (1단)",
        "슬라이드 부제목을 입력하세요",
        [[
            ("내용을 입력하세요", 0),
            ("내용을 입력하세요", 1),
            ("내용을 입력하세요", 2),
        ]],
        footnotes=["각주 예시: 출처나 참고 내용을 이렇게 입력하세요"],
    )

    add_content_slide(
        prs,
        "슬라이드 제목을 입력하세요 (2단)",
        "슬라이드 부제목을 입력하세요",
        [
            [("내용을 입력하세요", 0)] * 3,
            [("내용을 입력하세요", 0)] * 3,
        ],
    )

    add_content_slide(
        prs,
        "슬라이드 제목을 입력하세요 (3단)",
        "슬라이드 부제목을 입력하세요",
        [
            [("내용을 입력하세요", 0)] * 2,
            [("내용을 입력하세요", 0)] * 2,
            [("내용을 입력하세요", 0)] * 2,
        ],
    )

    add_content_slide(
        prs,
        "슬라이드 제목을 입력하세요 (4단)",
        "슬라이드 부제목을 입력하세요",
        [
            [("내용을 입력하세요", 0)] * 2,
            [("내용을 입력하세요", 0)] * 2,
            [("내용을 입력하세요", 0)] * 2,
            [("내용을 입력하세요", 0)] * 2,
        ],
    )

    add_image_slide(
        prs,
        "슬라이드 제목을 입력하세요 (이미지 1단)",
        "슬라이드 부제목을 입력하세요",
        [os.path.join(ASSETS, "bg.jpg")],
    )

    add_image_slide(
        prs,
        "슬라이드 제목을 입력하세요 (이미지 2단)",
        "슬라이드 부제목을 입력하세요",
        [os.path.join(ASSETS, "bg.jpg"), os.path.join(ASSETS, "bg1.jpg")],
    )

    add_image_slide(
        prs,
        "슬라이드 제목을 입력하세요 (이미지 3단)",
        "슬라이드 부제목을 입력하세요",
        [os.path.join(ASSETS, "bg.jpg"), os.path.join(ASSETS, "bg1.jpg"), os.path.join(ASSETS, "bg2.jpg")],
    )

    add_image_slide(
        prs,
        "슬라이드 제목을 입력하세요 (이미지 4단)",
        "슬라이드 부제목을 입력하세요",
        [
            os.path.join(ASSETS, "bg.jpg"),
            os.path.join(ASSETS, "bg1.jpg"),
            os.path.join(ASSETS, "bg2.jpg"),
            os.path.join(ASSETS, "logo_lab.png"),
        ],
    )

    add_table_slide(
        prs,
        "슬라이드 제목을 입력하세요 (표)",
        "슬라이드 부제목을 입력하세요",
        ["항목", "내용", "비고"],
        [
            ["항목 1", "내용을 입력하세요", "-"],
            ["항목 2", "내용을 입력하세요", "-"],
            ["항목 3", "내용을 입력하세요", "-"],
        ],
    )

    add_lead_slide(prs, "감사합니다", "", [])

    total = len(prs.slides)
    for i, slide in enumerate(prs.slides, start=1):
        add_page_number(slide, i, total)

    os.makedirs(os.path.dirname(OUT), exist_ok=True)
    prs.save(OUT)
    print(f"saved: {OUT}")


if __name__ == "__main__":
    build_demo()
